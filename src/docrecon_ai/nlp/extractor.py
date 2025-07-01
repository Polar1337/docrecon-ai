"""
Text extraction from various document formats

Extracts text content from different file types including:
- PDF documents
- Microsoft Office documents (Word, Excel, PowerPoint)
- Plain text files
- HTML files
- RTF documents
- OpenOffice documents
"""

import os
import re
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any, List
import logging

logger = logging.getLogger(__name__)

# Optional text extraction dependencies
try:
    import textract
    TEXTRACT_AVAILABLE = True
except ImportError:
    TEXTRACT_AVAILABLE = False
    logger.warning("textract not available. Install: pip install textract")

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF not available. Install: pip install PyMuPDF")

try:
    from docx import Document
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False
    logger.warning("python-docx not available. Install: pip install python-docx")

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not available. Install: pip install openpyxl")

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Install: pip install python-pptx")

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False
    logger.warning("BeautifulSoup not available. Install: pip install beautifulsoup4")


class TextExtractor:
    """
    Extracts text content from various document formats.
    
    Uses multiple extraction methods and libraries for best compatibility
    and text quality across different file types.
    """
    
    def __init__(self, config: Optional[Any] = None):
        """
        Initialize text extractor.
        
        Args:
            config: Configuration object with NLP settings
        """
        self.config = config
        self.logger = logging.getLogger(self.__class__.__name__)
        
        # Text processing settings
        self.max_text_length = getattr(config.nlp, 'max_text_length', 10000) if config else 10000
        self.clean_text = True
        
        # Statistics
        self.files_processed = 0
        self.extraction_errors = 0
    
    def extract_text(self, file_path: str, file_extension: str = None) -> Dict[str, Any]:
        """
        Extract text from a file.
        
        Args:
            file_path: Path to the file
            file_extension: File extension (auto-detected if not provided)
            
        Returns:
            dict: Extraction result with text content and metadata
        """
        if not file_extension:
            file_extension = Path(file_path).suffix.lower()
        
        result = {
            'text': '',
            'length': 0,
            'encoding': None,
            'extraction_method': None,
            'success': False,
            'error': None,
            'metadata': {}
        }
        
        try:
            # Select extraction method based on file type
            if file_extension == '.pdf':
                result = self._extract_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                result = self._extract_word(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                result = self._extract_excel(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                result = self._extract_powerpoint(file_path)
            elif file_extension in ['.txt', '.md']:
                result = self._extract_text_file(file_path)
            elif file_extension in ['.html', '.htm']:
                result = self._extract_html(file_path)
            elif file_extension == '.rtf':
                result = self._extract_rtf(file_path)
            elif file_extension in ['.odt', '.ods']:
                result = self._extract_openoffice(file_path)
            else:
                # Try generic textract as fallback
                result = self._extract_generic(file_path)
            
            # Clean and truncate text if needed
            if result['success'] and result['text']:
                result['text'] = self._clean_text(result['text'])
                result['length'] = len(result['text'])
                
                if self.max_text_length > 0 and result['length'] > self.max_text_length:
                    result['text'] = result['text'][:self.max_text_length]
                    result['length'] = self.max_text_length
                    result['metadata']['truncated'] = True
            
            self.files_processed += 1
            
        except Exception as e:
            self.logger.error(f"Error extracting text from {file_path}: {e}")
            result['error'] = str(e)
            self.extraction_errors += 1
        
        return result
    
    def _extract_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PDF using PyMuPDF"""
        result = {'text': '', 'success': False, 'extraction_method': 'pymupdf', 'metadata': {}}
        
        if not PYMUPDF_AVAILABLE:
            return self._extract_generic(file_path)
        
        try:
            doc = fitz.open(file_path)
            text_parts = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    text_parts.append(text)
            
            doc.close()
            
            result['text'] = '\n\n'.join(text_parts)
            result['success'] = True
            result['metadata']['pages'] = len(doc)
            
        except Exception as e:
            self.logger.error(f"PyMuPDF extraction failed for {file_path}: {e}")
            # Fallback to textract
            return self._extract_generic(file_path)
        
        return result
    
    def _extract_word(self, file_path: str) -> Dict[str, Any]:
        """Extract text from Word documents"""
        result = {'text': '', 'success': False, 'extraction_method': 'python-docx', 'metadata': {}}
        
        if not PYTHON_DOCX_AVAILABLE:
            return self._extract_generic(file_path)
        
        try:
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            result['text'] = '\n'.join(text_parts)
            result['success'] = True
            result['metadata']['paragraphs'] = len(doc.paragraphs)
            result['metadata']['tables'] = len(doc.tables)
            
        except Exception as e:
            self.logger.error(f"python-docx extraction failed for {file_path}: {e}")
            # Fallback to textract
            return self._extract_generic(file_path)
        
        return result
    
    def _extract_excel(self, file_path: str) -> Dict[str, Any]:
        """Extract text from Excel files"""
        result = {'text': '', 'success': False, 'extraction_method': 'openpyxl', 'metadata': {}}
        
        if not OPENPYXL_AVAILABLE:
            return self._extract_generic(file_path)
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())
                    if row_text:
                        sheet_text.append(' | '.join(row_text))
                
                if sheet_text:
                    text_parts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_text))
            
            result['text'] = '\n\n'.join(text_parts)
            result['success'] = True
            result['metadata']['sheets'] = len(workbook.sheetnames)
            
        except Exception as e:
            self.logger.error(f"openpyxl extraction failed for {file_path}: {e}")
            # Fallback to textract
            return self._extract_generic(file_path)
        
        return result
    
    def _extract_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """Extract text from PowerPoint files"""
        result = {'text': '', 'success': False, 'extraction_method': 'python-pptx', 'metadata': {}}
        
        if not PYTHON_PPTX_AVAILABLE:
            return self._extract_generic(file_path)
        
        try:
            presentation = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
            
            result['text'] = '\n\n'.join(text_parts)
            result['success'] = True
            result['metadata']['slides'] = len(presentation.slides)
            
        except Exception as e:
            self.logger.error(f"python-pptx extraction failed for {file_path}: {e}")
            # Fallback to textract
            return self._extract_generic(file_path)
        
        return result
    
    def _extract_text_file(self, file_path: str) -> Dict[str, Any]:
        """Extract text from plain text files"""
        result = {'text': '', 'success': False, 'extraction_method': 'direct', 'metadata': {}}
        
        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                
                result['text'] = text
                result['success'] = True
                result['encoding'] = encoding
                break
                
            except UnicodeDecodeError:
                continue
            except Exception as e:
                self.logger.error(f"Error reading text file {file_path}: {e}")
                break
        
        return result
    
    def _extract_html(self, file_path: str) -> Dict[str, Any]:
        """Extract text from HTML files"""
        result = {'text': '', 'success': False, 'extraction_method': 'beautifulsoup', 'metadata': {}}
        
        if not BS4_AVAILABLE:
            return self._extract_text_file(file_path)
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text()
            
            result['text'] = text
            result['success'] = True
            
        except Exception as e:
            self.logger.error(f"BeautifulSoup extraction failed for {file_path}: {e}")
            # Fallback to direct text extraction
            return self._extract_text_file(file_path)
        
        return result
    
    def _extract_rtf(self, file_path: str) -> Dict[str, Any]:
        """Extract text from RTF files"""
        # RTF extraction is complex, use textract as primary method
        return self._extract_generic(file_path)
    
    def _extract_openoffice(self, file_path: str) -> Dict[str, Any]:
        """Extract text from OpenOffice documents"""
        # OpenOffice extraction requires specialized libraries, use textract
        return self._extract_generic(file_path)
    
    def _extract_generic(self, file_path: str) -> Dict[str, Any]:
        """Generic text extraction using textract"""
        result = {'text': '', 'success': False, 'extraction_method': 'textract', 'metadata': {}}
        
        if not TEXTRACT_AVAILABLE:
            result['error'] = "No suitable text extraction library available"
            return result
        
        try:
            text = textract.process(file_path).decode('utf-8')
            result['text'] = text
            result['success'] = True
            
        except Exception as e:
            self.logger.error(f"textract extraction failed for {file_path}: {e}")
            result['error'] = str(e)
        
        return result
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not self.clean_text:
            return text
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x84\x86-\x9f]', '', text)
        
        # Normalize line breaks
        text = re.sub(r'\r\n|\r', '\n', text)
        
        # Remove excessive line breaks
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def get_supported_formats(self) -> List[str]:
        """
        Get list of supported file formats.
        
        Returns:
            list: Supported file extensions
        """
        formats = ['.txt', '.md']
        
        if PYMUPDF_AVAILABLE:
            formats.append('.pdf')
        
        if PYTHON_DOCX_AVAILABLE:
            formats.extend(['.docx', '.doc'])
        
        if OPENPYXL_AVAILABLE:
            formats.extend(['.xlsx', '.xls'])
        
        if PYTHON_PPTX_AVAILABLE:
            formats.extend(['.pptx', '.ppt'])
        
        if BS4_AVAILABLE:
            formats.extend(['.html', '.htm'])
        
        if TEXTRACT_AVAILABLE:
            formats.extend(['.rtf', '.odt', '.ods'])
        
        return sorted(list(set(formats)))
    
    def get_statistics(self) -> Dict[str, int]:
        """
        Get extraction statistics.
        
        Returns:
            dict: Statistics about text extraction
        """
        return {
            'files_processed': self.files_processed,
            'extraction_errors': self.extraction_errors,
        }
    
    def reset_statistics(self):
        """Reset extraction statistics"""
        self.files_processed = 0
        self.extraction_errors = 0

